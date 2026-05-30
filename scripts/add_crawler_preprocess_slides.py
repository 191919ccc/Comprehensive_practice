from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt


PPT = next(Path("docs/ppt").glob("股票实时流分析平台实时数据分析模块汇报.pptx"))
OUT = Path("docs/ppt/stock_realtime_with_crawler_preprocess.pptx")

prs = Presentation(str(PPT))

EMU_PER_INCH = 914400


def emu(value: float) -> int:
    return int(value * EMU_PER_INCH)


NAVY = RGBColor(15, 34, 56)
BLUE = RGBColor(32, 98, 167)
CYAN = RGBColor(40, 181, 194)
GREEN = RGBColor(62, 156, 106)
ORANGE = RGBColor(226, 139, 49)
LIGHT = RGBColor(247, 250, 253)
MID = RGBColor(100, 116, 139)
DARK = RGBColor(30, 41, 59)
BORDER = RGBColor(210, 220, 232)
WHITE = RGBColor(255, 255, 255)

layout = prs.slide_layouts[0]
new_slides = []


def make_slide(bg=LIGHT):
    slide = prs.slides.add_slide(layout)
    for shape in list(slide.shapes):
        shape.element.getparent().remove(shape.element)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg
    new_slides.append(slide)
    return slide


def add_textbox(slide, x, y, w, h, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = emu(0.04)
    tf.margin_right = emu(0.04)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.62, 0.42, 11.0, 0.42, title, 23, NAVY, True)
    if subtitle:
        add_textbox(slide, 0.65, 0.88, 10.6, 0.28, subtitle, 10.5, MID)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(0.62), emu(1.18), emu(1.15), emu(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = CYAN
    line.line.fill.background()


def add_card(slide, x, y, w, h, title, body, accent=BLUE):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, emu(x), emu(y), emu(w), emu(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = BORDER
    rect.line.width = Pt(1)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(x), emu(y), emu(0.08), emu(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_textbox(slide, x + 0.18, y + 0.15, w - 0.32, 0.28, title, 13, NAVY, True)
    tb = slide.shapes.add_textbox(emu(x + 0.18), emu(y + 0.52), emu(w - 0.32), emu(h - 0.62))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = emu(0.02)
    tf.margin_right = emu(0.02)
    for i, line in enumerate(body):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.space_after = Pt(4)
        run = p.runs[0]
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(9.2)
        run.font.color.rgb = DARK


def add_flow(slide, items, y=3.0):
    colors = [BLUE, CYAN, GREEN, ORANGE, RGBColor(112, 86, 190)]
    x0 = 0.75
    gap = 0.25
    boxw = (11.6 - gap * (len(items) - 1)) / len(items)
    for i, item in enumerate(items):
        x = x0 + i * (boxw + gap)
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, emu(x), emu(y), emu(boxw), emu(0.72))
        shp.fill.solid()
        shp.fill.fore_color.rgb = colors[i % len(colors)]
        shp.line.fill.background()
        tf = shp.text_frame
        tf.clear()
        tf.margin_left = emu(0.06)
        tf.margin_right = emu(0.06)
        tf.margin_top = emu(0.08)
        p = tf.paragraphs[0]
        p.text = item
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(10.5)
        run.font.bold = True
        run.font.color.rgb = WHITE
        if i < len(items) - 1:
            add_textbox(slide, x + boxw + 0.02, y + 0.2, 0.2, 0.25, "→", 17, MID, True, PP_ALIGN.CENTER)


def add_footer(slide, idx_text="数据爬虫与预处理"):
    add_textbox(slide, 0.65, 7.08, 4.5, 0.2, idx_text, 8.5, MID)
    add_textbox(slide, 10.7, 7.08, 1.7, 0.2, "股票实时流分析平台", 8.5, MID, False, PP_ALIGN.RIGHT)


s = make_slide(NAVY)
add_textbox(s, 0.82, 0.65, 7.9, 0.45, "数据爬虫与预处理模块", 27, WHITE, True)
add_textbox(s, 0.86, 1.22, 9.7, 0.35, "在实时数据分析前，先完成行情来源接入、字段解析、质量过滤和标准化输出", 13, RGBColor(210, 230, 245))
add_flow(s, ["行情源", "接口采集", "字段解析", "质量过滤", "标准事件"], 3.0)
add_textbox(s, 0.88, 5.08, 10.2, 0.6, "这一部分是实时链路的入口，决定后续 Kafka、Spark、MySQL/HDFS 中数据的完整性和可信度。", 15, RGBColor(235, 248, 255))

s = make_slide(LIGHT)
add_title(s, "数据来源与采集方式", "优先使用可直接请求的行情接口，避免 Selenium 带来的高延迟和不稳定")
add_card(s, 0.72, 1.55, 2.75, 2.05, "新浪财经", ["使用 requests 获取 A 股实时行情", "返回 JavaScript 字符串", "适合实时演示中的快速报价"], BLUE)
add_card(s, 3.75, 1.55, 2.75, 2.05, "腾讯财经", ["作为实时行情备用源", "返回 ~ 分隔字符串", "用于主源失败后的兜底"], CYAN)
add_card(s, 6.78, 1.55, 2.75, 2.05, "东方财富", ["通过 push2 API 获取行情", "返回 JSON 数据", "支持 A 股和港股报价"], GREEN)
add_card(s, 9.81, 1.55, 2.25, 2.05, "AKShare", ["获取历史日线数据", "用于冷启动和模型训练", "支持指数数据导入"], ORANGE)
add_textbox(s, 0.85, 4.22, 10.8, 0.28, "核心代码文件", 14, NAVY, True)
add_flow(s, ["stock_sources.py", "stock_producer.py", "cold_start.py", "stock_utils.py"], 4.75)
add_footer(s)

s = make_slide(RGBColor(250, 252, 255))
add_title(s, "爬虫采集流程设计", "Producer 按股票池循环采集，并通过多源兜底保证实时链路不断流")
add_flow(s, ["读取股票池", "多线程请求", "多源兜底", "行情校验", "发送 Kafka"], 1.72)
add_card(s, 0.78, 3.05, 3.3, 2.0, "并发采集", ["ThreadPoolExecutor 同时请求多只股票", "减少单轮行情采集耗时", "线程数由 STOCK_PRODUCER_MAX_WORKERS 控制"], BLUE)
add_card(s, 4.35, 3.05, 3.3, 2.0, "失败隔离", ["单只股票请求失败只跳过该股票", "不会中断整个 Producer 进程", "日志输出失败原因便于排查"], GREEN)
add_card(s, 7.92, 3.05, 3.3, 2.0, "Kafka 输出", ["股票代码作为消息 key", "行情事件序列化为 JSON", "Spark 后续持续消费实时 topic"], ORANGE)
add_textbox(s, 0.9, 5.82, 10.9, 0.42, "讲解重点：爬虫模块不是一次性下载文件，而是持续采集行情事件，为实时流计算提供不断更新的数据输入。", 13, DARK)
add_footer(s)

s = make_slide(LIGHT)
add_title(s, "字段解析与标准化输出", "不同来源返回格式不同，最终统一为 Kafka 和 Spark 可消费的标准行情事件")
add_card(s, 0.78, 1.55, 3.45, 2.1, "返回格式差异", ["新浪：JavaScript 字符串", "腾讯：~ 分隔字段", "东方财富：JSON 字段", "AKShare：DataFrame 行记录"], BLUE)
add_card(s, 4.55, 1.55, 3.45, 2.1, "统一解析方法", ["正则提取字符串 payload", "response.json() 读取接口 JSON", "safe_float / safe_int 安全转换", "detect_market 识别市场"], CYAN)
add_card(s, 8.32, 1.55, 3.45, 2.1, "标准事件字段", ["symbol、company_name、market", "open/high/low/last_price", "change_pct、volume、turnover", "event_time、source、event_id"], GREEN)
add_textbox(s, 0.85, 4.25, 11.0, 0.28, "标准化后的事件结构", 14, NAVY, True)
add_flow(s, ["股票标识", "价格字段", "成交字段", "时间来源", "唯一事件 ID"], 4.78)
add_textbox(s, 0.88, 6.05, 10.6, 0.35, "统一结构的好处：后续 Kafka 传输、Spark 解析、MySQL 写库和前端展示都不需要关心原始数据源差异。", 12.5, DARK)
add_footer(s)

s = make_slide(RGBColor(249, 251, 253))
add_title(s, "数据预处理与质量控制", "在实时计算前先处理缺失值、异常值、重复数据和训练特征")
add_card(s, 0.72, 1.48, 2.75, 2.18, "缺失值处理", ["空值、--、非法字符串转默认值", "价格和成交量统一转数值", "滚动指标产生的 NaN 统一填充"], BLUE)
add_card(s, 3.75, 1.48, 2.75, 2.18, "异常值过滤", ["last_price > 0", "volume > 0", "A 股涨跌幅限制更严格", "港股允许更大波动范围"], ORANGE)
add_card(s, 6.78, 1.48, 2.75, 2.18, "去重与压缩", ["按 event_id 去重", "过滤未变化行情快照", "训练集压缩连续不变价格"], CYAN)
add_card(s, 9.81, 1.48, 2.25, 2.18, "特征构造", ["收益率、动量、均线", "成交量变化、波动率", "行业和指数相对特征"], GREEN)
add_flow(s, ["原始行情", "安全转换", "质量过滤", "去重压缩", "可分析数据"], 4.55)
add_textbox(s, 0.86, 5.82, 10.9, 0.44, "这一阶段保证进入实时分析和机器学习模块的数据是结构化、可信、可比较的行情数据。", 13, DARK)
add_footer(s)

sld_id_list = prs.slides._sldIdLst
ids = list(sld_id_list)
new_ids = ids[-len(new_slides) :]
for slide_id in new_ids:
    sld_id_list.remove(slide_id)
for offset, slide_id in enumerate(new_ids):
    sld_id_list.insert(1 + offset, slide_id)

prs.save(str(OUT))
print(OUT)
print(f"slides={len(prs.slides)}")
